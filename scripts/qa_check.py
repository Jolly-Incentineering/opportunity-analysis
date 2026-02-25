"""
qa_check.py — Run QA checks on the Excel model and vF PowerPoint deck.

Usage:  python .claude/scripts/qa_check.py --company "Company Name"

Checks:
    Excel:  formula counts, comment coverage, value reasonableness
    PPT:    red text, placeholders, dollar formatting, banner fill, uppercase K
    Cross:  key values approximately match between Excel and PPT
"""
import sys, re, argparse

sys.stdout.reconfigure(encoding="utf-8", errors="replace")

try:
    from openpyxl import load_workbook
except ImportError:
    print("ERROR: openpyxl not installed. Run: pip install openpyxl"); sys.exit(1)
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx"); sys.exit(1)

from jolly_utils import CLIENTS_DIR, FORMULA_COUNTS, count_formulas

PASS, FAIL, WARN = "[PASS]", "[FAIL]", "[WARN]"
RED = RGBColor(0xFF, 0x00, 0x00)
PLACEHOLDER_RE = re.compile(r"\[[ \w]*\]")
RAW_DOLLAR_RE = re.compile(r"\$[\d,]{5,}")
UPPERCASE_K_RE = re.compile(r"\$\d+K\b")


def find_file(company: str, subfolder: str, pattern: str) -> str:
    """Find the most-recent file matching *pattern* inside the client folder."""
    folder = CLIENTS_DIR / company / subfolder
    matches = sorted(folder.glob(pattern))
    if not matches:
        raise FileNotFoundError(f"No file matching '{pattern}' in {folder}")
    return str(matches[-1])


def find_vf_deck(company: str) -> str:
    """Locate the vF PowerPoint deck (case-insensitive)."""
    folder = CLIENTS_DIR / company / "2. Presentations"
    for pat in ("*vF*.pptx", "*vf*.pptx"):
        matches = sorted(folder.glob(pat))
        if matches:
            return str(matches[-1])
    raise FileNotFoundError(f"No vF deck in {folder}")


def detect_industry(wb) -> str:
    if "Campaigns" not in wb.sheetnames:
        return "qsr"
    return "manufacturing" if count_formulas(wb["Campaigns"]) > 200 else "qsr"


def check_excel(company: str) -> dict:
    results = {}
    print("\n=== EXCEL MODEL ===")
    try:
        model_path = find_file(company, "1. Model", "*.xlsx")
        print(f"File: {model_path}")
    except FileNotFoundError as e:
        print(f"  {FAIL} Model file not found: {e}"); return {}

    wb_formulas = load_workbook(model_path, data_only=False)
    wb_values = load_workbook(model_path, data_only=True)
    industry = detect_industry(wb_formulas)
    expected = FORMULA_COUNTS.get(industry, {})
    print(f"  Detected industry: {industry}")

    # 1. Formula counts
    for sheet, exp in expected.items():
        if sheet not in wb_formulas.sheetnames:
            print(f"  {WARN} {sheet} sheet not found"); results[f"{sheet}_formulas"] = None; continue
        actual = count_formulas(wb_formulas[sheet])
        ok = actual == exp
        suffix = "" if ok else " -- formulas may be overwritten"
        print(f"  {PASS if ok else FAIL} {sheet} formulas: {actual} (expected {exp}){suffix}")
        results[f"{sheet}_formulas"] = ok

    # 2. Comment coverage on Inputs sheet
    if "Inputs" in wb_formulas.sheetnames:
        ws = wb_formulas["Inputs"]
        missing = [
            cell.coordinate
            for row in ws.iter_rows(min_row=1, max_row=100, min_col=3, max_col=5)
            for cell in row
            if cell.value is not None
            and not (isinstance(cell.value, str) and cell.value.startswith("="))
            and cell.comment is None
        ]
        if not missing:
            print(f"  {PASS} Comment coverage: all value cells commented"); results["comments"] = True
        else:
            print(f"  {FAIL} Comment coverage: {len(missing)} cells missing comments: {missing[:10]}")
            results["comments"] = False
    else:
        print(f"  {WARN} Inputs sheet not found"); results["comments"] = None

    # 3. Value reasonableness
    if "Inputs" in wb_values.sheetnames:
        ws = wb_values["Inputs"]
        issues = []
        for row in ws.iter_rows(min_row=1, max_row=100, min_col=3, max_col=5):
            for cell in row:
                v = cell.value
                if v is None or not isinstance(v, (int, float)):
                    continue
                if v < 0:
                    issues.append(f"{cell.coordinate}={v} (negative)")
                if isinstance(v, float) and v != int(v) and v > 10:
                    issues.append(f"{cell.coordinate}={v} (fractional headcount?)")
        if not issues:
            print(f"  {PASS} Value reasonableness: no obvious issues"); results["reasonableness"] = True
        else:
            print(f"  {WARN} Value reasonableness: check these cells: {issues[:5]}")
            results["reasonableness"] = None
    return results


def check_ppt(company: str) -> dict:
    """Single-pass PPT check: red text, placeholders, dollar fmt, uppercase K, banner."""
    results = {}
    print("\n=== POWERPOINT DECK ===")
    try:
        vf_path = find_vf_deck(company)
        print(f"File: {vf_path}")
    except FileNotFoundError as e:
        print(f"  {FAIL} vF deck not found: {e}"); return {}

    prs = Presentation(vf_path)
    red_runs, placeholders, raw_dollars, uppercase_k = [], [], [], []
    banner_ok = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if PLACEHOLDER_RE.search(text):
                placeholders.append(text[:60])
            raw_dollars.extend(RAW_DOLLAR_RE.findall(text))
            uppercase_k.extend(UPPERCASE_K_RE.findall(text))
            if not banner_ok and ("MM" in text or "k" in text) and "quantified" in text:
                if "$[ ]" not in text and "[ ] quantified" not in text:
                    banner_ok = True
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.rgb == RED:
                            red_runs.append(run.text[:40])
                    except Exception:
                        pass

    # Report each check
    for key, items, pass_msg, fail_fmt in [
        ("red_text", red_runs, "No red text",
         "{n} red text runs found: {s}"),
        ("placeholders", placeholders, "No unfilled placeholders",
         "{n} unfilled placeholders: {s}"),
        ("dollar_format", raw_dollars, "Dollar formatting: all amounts use K/MM",
         "{n} raw dollar amounts: {s}"),
        ("lowercase_k", uppercase_k, "No uppercase $K (all lowercase $k)",
         "Uppercase $K found: {s} (should be lowercase)"),
    ]:
        if not items:
            print(f"  {PASS} {pass_msg}"); results[key] = True
        else:
            print(f"  {FAIL} {fail_fmt.format(n=len(items), s=items[:5])}")
            results[key] = False

    if banner_ok:
        print(f"  {PASS} Banner appears filled"); results["banner"] = True
    else:
        print(f"  {WARN} Banner may not be filled - check summary slide manually")
        results["banner"] = None
    return results


def check_cross_validation(company: str) -> dict:
    results = {}
    print("\n=== CROSS-VALIDATION (Excel vs PPT) ===")

    excel_values = {}
    try:
        model_path = find_file(company, "1. Model", "*.xlsx")
        wb = load_workbook(model_path, data_only=True)
        if "Inputs" in wb.sheetnames:
            ws = wb["Inputs"]
            for r in range(5, 10):
                label, val = ws.cell(row=r, column=2).value, ws.cell(row=r, column=3).value
                if val is not None and label:
                    excel_values[str(label).strip()] = val
    except Exception as e:
        print(f"  {WARN} Could not read Excel: {e}"); return {}

    if not excel_values:
        print(f"  {WARN} No values extracted from Excel Inputs sheet"); return {}

    try:
        prs = Presentation(find_vf_deck(company))
        ppt_text = " ".join(
            shape.text_frame.text for slide in prs.slides
            for shape in slide.shapes if shape.has_text_frame
        )
    except Exception as e:
        print(f"  {WARN} Could not read PPT: {e}"); return {}

    mismatches = 0
    for label, value in list(excel_values.items())[:5]:
        if not isinstance(value, (int, float)) or value == 0:
            continue
        val_str = str(int(value))
        val_mm = f"{value / 1_000_000:.0f}" if value >= 1_000_000 else None
        found = val_str in ppt_text or (val_mm and val_mm in ppt_text)
        detail = "found in deck" if found else "NOT found in deck (may be formatted differently)"
        print(f"  {PASS if found else WARN} {label}: {value} {detail}")
        if not found:
            mismatches += 1
    results["cross_validation"] = mismatches == 0
    return results


def main():
    parser = argparse.ArgumentParser(description="QA check for intro deck package")
    parser.add_argument("--company", required=True, help="Company name (must match Clients/ folder)")
    args = parser.parse_args()
    company = args.company

    print(f"\n=== qa_check.py | {company} ===")
    all_results = {}
    all_results.update(check_excel(company))
    all_results.update(check_ppt(company))
    all_results.update(check_cross_validation(company))

    failures = [k for k, v in all_results.items() if v is False]
    warnings = [k for k, v in all_results.items() if v is None]
    total = len(all_results)
    passed = sum(1 for v in all_results.values() if v is True)

    print(f"\n=== SUMMARY ===\n  {passed}/{total} checks passed")
    if failures:
        print(f"\n  FAILURES ({len(failures)}):")
        for f in failures:
            print(f"    - {f}")
    if warnings:
        print(f"\n  WARNINGS ({len(warnings)}) - verify manually:")
        for w in warnings:
            print(f"    - {w}")
    if not failures:
        print(f"\n  OVERALL: PASS - ready for delivery")
    else:
        print(f"\n  OVERALL: FAIL - fix issues above and re-run")
        print(f'\n  Re-run: python .claude/scripts/qa_check.py --company "{company}"')


if __name__ == "__main__":
    main()
