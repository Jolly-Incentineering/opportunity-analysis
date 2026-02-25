#!/usr/bin/env python3
"""
deck_engine.py — Consolidated deck operations for the Jolly workflow.
=====================================================================
Single CLI tool replacing inline Python in skill prompts.

Usage:
    python3 deck_engine.py <action> [options]

Actions:
    fill-banners      Replace bracket placeholders with values from research JSON
    format-dollars    Reformat raw dollar amounts ($1234567 → $1.2MM)
    find-placeholders Find all remaining [...] template tokens
    set-title         Set document title on .pptx or .xlsx
    set-pdf-title     Set PDF /Title metadata from the source presentation
    finalize          Convert red Macabacus text to black/white for delivery
    copy-vf           Copy master deck to vF delivery copy

Requirements: python-pptx, openpyxl, pypdf
"""
import argparse
import json
import os
import re
import shutil
import sys
from pathlib import Path

# ---------------------------------------------------------------------------
# Graceful dependency checks
# ---------------------------------------------------------------------------
def _require(module_name, pip_name=None):
    """Import a module or exit with a clear install message."""
    try:
        return __import__(module_name)
    except ImportError:
        pip_name = pip_name or module_name
        print(f"ERROR: {module_name} is required. Install with: pip install {pip_name}", file=sys.stderr)
        sys.exit(1)


# ---------------------------------------------------------------------------
# Dollar formatting — THE single standard
# ---------------------------------------------------------------------------
DOLLAR_FORMAT_HELP = """
Dollar formatting standard:
    $1M+       → $X.XMM   (1 decimal, uppercase MM, e.g. $21.6MM, $2.0MM)
    $1K–$999K  → $XXXk    (integer, lowercase k, e.g. $516k, $2k)
    Under $1K  → $XXX     (plain integer, no suffix)
"""

def format_dollars(value):
    """Format a numeric value per Jolly dollar standards.

    Returns formatted string (e.g. "$21.6MM", "$516k", "$850").
    """
    if value is None:
        return ""
    value = abs(float(value))

    if value >= 1_000_000:
        mm = value / 1_000_000
        return f"${mm:.1f}MM"
    elif value >= 1_000:
        k = int(round(value / 1_000))
        return f"${k}k"
    else:
        return f"${int(value)}"


# ---------------------------------------------------------------------------
# Macabacus detection
# ---------------------------------------------------------------------------
def is_macabacus_linked(run):
    """Return True if a text run has red font (Macabacus live link).

    Macabacus-linked runs have red text (R>200, G<100, B<100) and must
    never be edited — Macabacus refresh populates their values.
    """
    try:
        if run.font.color and run.font.color.type is not None:
            rgb = run.font.color.rgb
            if rgb and len(str(rgb)) >= 6:
                h = str(rgb)
                r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                return r > 200 and g < 100 and b < 100
    except Exception:
        pass
    return False


# ---------------------------------------------------------------------------
# Bracket placeholder regex — matches any [...] token
# ---------------------------------------------------------------------------
BRACKET_RE = re.compile(r'\[.*?\]')

# Raw dollar regex — $XXXXX+ with optional commas and decimals
RAW_DOLLAR_RE = re.compile(r'\$[\d,]{5,}(?:\.\d+)?')


# ---------------------------------------------------------------------------
# Actions
# ---------------------------------------------------------------------------
def action_fill_banners(args):
    """Fill bracket placeholders in a deck from research JSON.

    Reads campaign_details from research JSON and replaces banner
    placeholders with formatted EBITDA values.
    """
    pptx_mod = _require("pptx", "python-pptx")
    abs_path = os.path.abspath(args.file)
    prs = pptx_mod.Presentation(abs_path)

    # Load research data
    with open(args.research, "r", encoding="utf-8") as f:
        research = json.load(f)

    campaign_details = research.get("campaign_details", {})
    campaigns_selected = research.get("campaigns_selected", [])

    # Build replacement values
    # Total EBITDA = sum of all campaign ebitda_uplift_base
    total_ebitda = 0
    campaign_count = 0
    for camp in campaigns_selected:
        name = camp.get("name", camp) if isinstance(camp, dict) else camp
        detail = campaign_details.get(name, {})
        ebitda = detail.get("ebitda_uplift_base", 0)
        if ebitda:
            total_ebitda += ebitda
            campaign_count += 1

    total_formatted = format_dollars(total_ebitda).replace("$", "")  # e.g. "65.4MM"

    replacements_made = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                # Skip Macabacus-linked paragraphs
                if any(is_macabacus_linked(r) for r in para.runs):
                    continue

                full_text = "".join(run.text for run in para.runs)
                if not BRACKET_RE.search(full_text):
                    continue

                # Replace all bracket placeholders in the paragraph
                new_text = full_text

                # EBITDA placeholder: $[ ] or $[EBITDA] followed by MM
                new_text = re.sub(r'\$\[.*?\]\s*MM', f'${total_formatted}', new_text)
                # Standalone $[...] (EBITDA without MM suffix)
                new_text = re.sub(r'\$\[.*?\]', f'${total_formatted}', new_text)
                # Campaign count: [ ] quantified or [N] quantified
                new_text = re.sub(r'\[.*?\]\s*quantified', f'{campaign_count} quantified', new_text)
                # Generic remaining [...] — replace with campaign count
                new_text = BRACKET_RE.sub(str(campaign_count), new_text)

                if new_text != full_text:
                    # Write back: put all text in first run, clear others
                    if para.runs:
                        para.runs[0].text = new_text
                        for run in para.runs[1:]:
                            run.text = ""
                        replacements_made += 1

    prs.save(abs_path)
    result = {
        "replacements": replacements_made,
        "total_ebitda": total_ebitda,
        "total_formatted": f"${total_formatted}",
        "campaign_count": campaign_count,
    }
    print(json.dumps(result, indent=2))
    return result


def action_format_dollars(args):
    """Reformat raw dollar amounts in a deck.

    Finds patterns like $1234567 or $1,234,567 and reformats per standard.
    Skips Macabacus-linked runs. Optionally skips specific slides.
    """
    pptx_mod = _require("pptx", "python-pptx")
    abs_path = os.path.abspath(args.file)
    prs = pptx_mod.Presentation(abs_path)

    skip_slides = set()
    if args.skip_slides:
        skip_slides = {int(s) for s in args.skip_slides.split(",")}

    replacements_made = 0
    details = []

    for i, slide in enumerate(prs.slides, 1):
        if i in skip_slides:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if is_macabacus_linked(run):
                        continue
                    matches = RAW_DOLLAR_RE.findall(run.text)
                    if not matches:
                        continue
                    new_text = run.text
                    for m in matches:
                        raw_num = float(m.replace("$", "").replace(",", ""))
                        formatted = format_dollars(raw_num)
                        new_text = new_text.replace(m, formatted, 1)
                    if new_text != run.text:
                        old = run.text
                        run.text = new_text
                        replacements_made += 1
                        details.append({"slide": i, "old": old, "new": new_text})

    prs.save(abs_path)
    result = {"replacements": replacements_made, "details": details}
    print(json.dumps(result, indent=2))
    return result


def action_find_placeholders(args):
    """Find all remaining bracket placeholders in a deck."""
    pptx_mod = _require("pptx", "python-pptx")
    prs = pptx_mod.Presentation(os.path.abspath(args.file))

    results = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                for m in BRACKET_RE.finditer(text):
                    results.append({
                        "slide": i,
                        "shape": shape.name,
                        "match": m.group(),
                        "context": text[:120],
                    })

    print(json.dumps(results, indent=2))
    return results


def action_set_title(args):
    """Set document title on a .pptx or .xlsx file."""
    abs_path = os.path.abspath(args.file)
    ext = Path(abs_path).suffix.lower()

    if ext == ".pptx":
        pptx_mod = _require("pptx", "python-pptx")
        prs = pptx_mod.Presentation(abs_path)
        prs.core_properties.title = args.title
        prs.save(abs_path)
    elif ext == ".xlsx":
        openpyxl = _require("openpyxl")
        wb = openpyxl.load_workbook(abs_path)
        wb.properties.title = args.title
        wb.save(abs_path)
    else:
        print(f"ERROR: Unsupported file type: {ext}", file=sys.stderr)
        sys.exit(1)

    print(json.dumps({"file": abs_path, "title": args.title}))


def action_set_pdf_title(args):
    """Set PDF /Title metadata to match the source presentation's title."""
    pypdf = _require("pypdf")
    pptx_mod = _require("pptx", "python-pptx")

    abs_pdf = os.path.abspath(args.file)
    abs_pptx = os.path.abspath(args.from_pptx)

    # Read title from the presentation
    prs = pptx_mod.Presentation(abs_pptx)
    pdf_title = prs.core_properties.title
    if not pdf_title:
        pdf_title = Path(abs_pptx).stem  # fallback to filename without extension

    writer = pypdf.PdfWriter(clone_from=abs_pdf)
    writer.add_metadata({"/Title": pdf_title})
    writer.write(abs_pdf)

    print(json.dumps({"pdf": abs_pdf, "title": pdf_title}))


def action_finalize(args):
    """Convert red Macabacus text to black or white for client delivery.

    Dark backgrounds → white text. Light/no backgrounds → black text.
    """
    pptx_mod = _require("pptx", "python-pptx")
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_FILL_TYPE

    abs_path = os.path.abspath(args.file)
    prs = pptx_mod.Presentation(abs_path)
    fixed = 0

    def _bg_is_dark(shape):
        try:
            if hasattr(shape, "fill") and shape.fill.type == MSO_FILL_TYPE.SOLID:
                rgb_hex = str(shape.fill.fore_color.rgb)
                if len(rgb_hex) >= 6:
                    r, g, b = int(rgb_hex[0:2], 16), int(rgb_hex[2:4], 16), int(rgb_hex[4:6], 16)
                    return (r + g + b) / 3 < 100
        except Exception:
            pass
        return False

    for slide in prs.slides:
        for shape in slide.shapes:
            frames = []
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                frames.append((shape.text_frame, shape))
            if shape.shape_type == 19:  # Table
                try:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            frames.append((cell.text_frame, shape))
                except Exception:
                    pass

            dark = _bg_is_dark(shape)
            target_color = RGBColor(255, 255, 255) if dark else RGBColor(0, 0, 0)

            for tf, _ in frames:
                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.text.strip() and is_macabacus_linked(run):
                            run.font.color.rgb = target_color
                            fixed += 1

    prs.save(abs_path)
    print(json.dumps({"red_text_fixed": fixed}))


def action_copy_vf(args):
    """Copy master deck to vF delivery copy and update its title."""
    pptx_mod = _require("pptx", "python-pptx")

    src = os.path.abspath(args.src)
    dest = os.path.abspath(args.dest)

    shutil.copy2(src, dest)

    # Update title to match the vF filename
    prs = pptx_mod.Presentation(dest)
    prs.core_properties.title = Path(dest).stem
    prs.save(dest)

    print(json.dumps({"src": src, "dest": dest, "title": Path(dest).stem}))


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def main():
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

    parser = argparse.ArgumentParser(
        description="Jolly deck engine — consolidated deck operations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=DOLLAR_FORMAT_HELP,
    )
    sub = parser.add_subparsers(dest="action", required=True)

    # fill-banners
    p = sub.add_parser("fill-banners", help="Fill bracket placeholders from research JSON")
    p.add_argument("--file", required=True, help="Path to .pptx file")
    p.add_argument("--research", required=True, help="Path to research_output JSON")

    # format-dollars
    p = sub.add_parser("format-dollars", help="Reformat raw dollar amounts")
    p.add_argument("--file", required=True, help="Path to .pptx file")
    p.add_argument("--skip-slides", default=None, help="Comma-separated slide numbers to skip")

    # find-placeholders
    p = sub.add_parser("find-placeholders", help="Find remaining [...] tokens")
    p.add_argument("--file", required=True, help="Path to .pptx file")

    # set-title
    p = sub.add_parser("set-title", help="Set document title on .pptx or .xlsx")
    p.add_argument("--file", required=True, help="Path to .pptx or .xlsx file")
    p.add_argument("--title", required=True, help="Title to set")

    # set-pdf-title
    p = sub.add_parser("set-pdf-title", help="Set PDF title from source presentation")
    p.add_argument("--file", required=True, help="Path to .pdf file")
    p.add_argument("--from-pptx", required=True, help="Path to source .pptx")

    # finalize
    p = sub.add_parser("finalize", help="Convert red Macabacus text to black/white")
    p.add_argument("--file", required=True, help="Path to .pptx file")

    # copy-vf
    p = sub.add_parser("copy-vf", help="Copy master to vF delivery copy")
    p.add_argument("--src", required=True, help="Path to master .pptx")
    p.add_argument("--dest", required=True, help="Path for vF .pptx")

    args = parser.parse_args()

    actions = {
        "fill-banners": action_fill_banners,
        "format-dollars": action_format_dollars,
        "find-placeholders": action_find_placeholders,
        "set-title": action_set_title,
        "set-pdf-title": action_set_pdf_title,
        "finalize": action_finalize,
        "copy-vf": action_copy_vf,
    }

    actions[args.action](args)


if __name__ == "__main__":
    main()
