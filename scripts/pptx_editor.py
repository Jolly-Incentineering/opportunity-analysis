"""
PowerPoint Editor Agent
=======================
Creates and edits Jolly intro deck presentations.

Capabilities:
  - Copy template and rename for company
  - Fill Slide 8 EBITDA placeholder ("$[ ]  MM")
  - Update company name across all slides
  - Insert company logo into branded positions
  - Sync numbers from Excel model
  - Set document properties (title, category)
"""
import shutil
from pathlib import Path
from datetime import date

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from .jolly_utils import (
    BASE_DIR, get_template_paths, resolve_client_path, COMMENT_AUTHOR,
)


def _require_pptx():
    if not HAS_PPTX:
        raise ImportError(
            "python-pptx is required. Install with: pip install python-pptx"
        )


class PptxEditor:
    """Creates and edits Jolly PowerPoint intro decks."""

    def __init__(self, template_type="qsr"):
        _require_pptx()
        self.template_type = template_type.lower()
        _, self.ppt_template = get_template_paths(self.template_type)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------
    def create_deck(self, company_name, folder=None, date_str=None):
        """Copy PPT template to client folder. Returns Path."""
        client_dir = resolve_client_path(company_name, folder)
        ppt_dir = client_dir / "2. Presentations"
        ppt_dir.mkdir(parents=True, exist_ok=True)

        if date_str is None:
            date_str = date.today().strftime("%Y.%m.%d")

        filename = f"{company_name} Intro Deck ({date_str}).pptx"
        dest = ppt_dir / filename
        shutil.copy2(self.ppt_template, dest)
        return dest

    def fill_ebitda_placeholder(self, ppt_path, ebitda_mm, campaign_count):
        """Fill the Slide 8 'Rectangle 80' EBITDA placeholder.

        Template text:
          "Identified $[ ]  MM of EBITDA in highly actionable optimizations
           with just [ ] quantified Jolly incentive campaigns"

        Replaces $[ ] with dollar amount and [ ] with campaign count.
        """
        prs = Presentation(str(ppt_path))
        filled = False

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.name == "Rectangle 80" and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        full_text = "".join(run.text for run in para.runs)
                        if "$[ ]" in full_text:
                            new_text = full_text.replace(
                                "$[ ]", f"${ebitda_mm}"
                            ).replace(
                                "[ ]", str(campaign_count)
                            )
                            # Preserve formatting of first run
                            if para.runs:
                                for i, run in enumerate(para.runs):
                                    if i == 0:
                                        run.text = new_text
                                    else:
                                        run.text = ""
                            filled = True

        if filled:
            prs.save(str(ppt_path))
        return {"filled": filled, "ebitda_mm": ebitda_mm,
                "campaigns": campaign_count}

    def update_company_name(self, ppt_path, old_name, new_name):
        """Replace all occurrences of old_name with new_name in text."""
        prs = Presentation(str(ppt_path))
        replacements = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if old_name in run.text:
                                run.text = run.text.replace(
                                    old_name, new_name
                                )
                                replacements += 1

        if replacements > 0:
            prs.save(str(ppt_path))
        return {"replacements": replacements}

    def replace_text(self, ppt_path, replacements_dict):
        """Replace multiple text strings in the deck.

        *replacements_dict*: {"old_text": "new_text", ...}
        """
        prs = Presentation(str(ppt_path))
        counts = {k: 0 for k in replacements_dict}

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            for old, new in replacements_dict.items():
                                if old in run.text:
                                    run.text = run.text.replace(old, new)
                                    counts[old] += 1

        if any(v > 0 for v in counts.values()):
            prs.save(str(ppt_path))
        return {"replacements": counts}

    def insert_logo(self, ppt_path, logo_path, slide_index=0,
                    left_inches=0.5, top_inches=0.5,
                    width_inches=2.0):
        """Insert a logo image into a specific slide."""
        prs = Presentation(str(ppt_path))
        slide = prs.slides[slide_index]
        slide.shapes.add_picture(
            str(logo_path),
            Inches(left_inches), Inches(top_inches),
            width=Inches(width_inches),
        )
        prs.save(str(ppt_path))
        return {"slide": slide_index, "logo": str(logo_path)}

    def set_document_properties(self, ppt_path, title=None, category=None):
        """Set document metadata properties."""
        prs = Presentation(str(ppt_path))
        if title:
            prs.core_properties.title = title
        if category:
            prs.core_properties.category = category
        prs.save(str(ppt_path))
        return {"title": title, "category": category}

    def read_all_text(self, ppt_path):
        """Extract all text from deck for validation purposes.

        Returns dict with slide_number -> list of (shape_name, text).
        """
        prs = Presentation(str(ppt_path))
        result = {}

        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = "\n".join(
                        p.text for p in shape.text_frame.paragraphs
                    )
                    if text.strip():
                        slide_texts.append((shape.name, text.strip()))
            if slide_texts:
                result[i] = slide_texts

        return result

    def find_unfilled_placeholders(self, ppt_path):
        """Find any remaining template placeholders in the deck.

        Looks for patterns: $[ ], [ ], [Company Name], etc.
        """
        import re
        prs = Presentation(str(ppt_path))
        placeholders = []
        pattern = re.compile(r'\[[\s\w]*\]')

        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        matches = pattern.findall(text)
                        for m in matches:
                            placeholders.append({
                                "slide": i,
                                "shape": shape.name,
                                "placeholder": m,
                                "context": text[:100],
                            })

        return placeholders

    def sync_from_model(self, ppt_path, model_summary):
        """Sync key numbers from an Excel model summary into the deck.

        *model_summary*: dict from ExcelEditor.read_model_summary()
        """
        replacements = {}
        name = model_summary.get("company_name")
        revenue = model_summary.get("revenue")

        if revenue:
            revenue_mm = round(revenue / 1_000_000, 1)
            # Common patterns to replace
            replacements["$[ ]"] = f"${revenue_mm}"

        result = {"synced_fields": []}
        if replacements:
            self.replace_text(ppt_path, replacements)
            result["synced_fields"] = list(replacements.keys())

        return result

    def finalize(self, ppt_path):
        """Finalize presentation for client delivery.

        Converts all red text (Macabacus links) to black or white based on background:
        - Dark backgrounds → white text
        - Light/no backgrounds → black text

        Returns dict with count of text runs fixed.
        """
        from pptx.dml.color import RGBColor
        from pptx.enum.dml import MSO_FILL_TYPE

        ppt_path = Path(ppt_path)
        prs = Presentation(ppt_path)
        fixed_count = 0

        def get_background_type(shape):
            """Determine if shape has dark or light background"""
            try:
                if hasattr(shape, 'fill') and shape.fill.type == MSO_FILL_TYPE.SOLID:
                    if hasattr(shape.fill.fore_color, 'rgb'):
                        rgb_hex = str(shape.fill.fore_color.rgb)
                        if len(rgb_hex) == 6:
                            r = int(rgb_hex[0:2], 16)
                            g = int(rgb_hex[2:4], 16)
                            b = int(rgb_hex[4:6], 16)
                            avg = (r + g + b) / 3
                            # Dark if average RGB < 100 (navy, dark blue, black)
                            return 'DARK' if avg < 100 else 'LIGHT'
            except:
                pass
            return 'LIGHT'  # Default to light (black text)

        def fix_red_text_in_frame(text_frame, shape):
            """Find and fix red text in a text frame"""
            nonlocal fixed_count
            bg_type = get_background_type(shape)

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        try:
                            if run.font.color and run.font.color.type == 1:  # RGB
                                rgb_hex = str(run.font.color.rgb)
                                if len(rgb_hex) == 6:
                                    r = int(rgb_hex[0:2], 16)
                                    g = int(rgb_hex[2:4], 16)
                                    b = int(rgb_hex[4:6], 16)

                                    # RED TEXT: R>200, G<100, B<100
                                    if r > 200 and g < 100 and b < 100:
                                        if bg_type == 'DARK':
                                            run.font.color.rgb = RGBColor(255, 255, 255)
                                        else:
                                            run.font.color.rgb = RGBColor(0, 0, 0)
                                        fixed_count += 1
                        except:
                            pass

        # Process all slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    fix_red_text_in_frame(shape.text_frame, shape)

                # Tables
                if shape.shape_type == 19:
                    try:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                fix_red_text_in_frame(cell.text_frame, shape)
                    except:
                        pass

        # Save
        prs.save(ppt_path)

        return {"red_text_fixed": fixed_count}
